import base64
import json
import sys
import tempfile
import unittest
import zipfile
from pathlib import Path

import openpyxl
import zstandard
from pptx import Presentation
from pptx.util import Inches

ROOT = Path(__file__).resolve().parents[1]
if str(ROOT) not in sys.path:
    sys.path.insert(0, str(ROOT))

from converted_document import DocumentConversionError
from converters.pptx import PptxConverter, PresentationLimits
from converters.spreadsheet import (
    CsvConverter,
    OoxmlZipLimits,
    SpreadsheetLimits,
    XlsxConverter,
)
from source_map import strip_source_unit_markers


def read_source_map(path: Path) -> list[dict]:
    with (
        path.open("rb") as raw,
        zstandard.ZstdDecompressor().stream_reader(raw) as reader,
    ):
        payload = reader.read().decode("utf-8")
    return [json.loads(line) for line in payload.splitlines()]


class PresentationSpreadsheetConverterTests(unittest.TestCase):
    def setUp(self):
        self.temporary_directory = tempfile.TemporaryDirectory()
        self.root = Path(self.temporary_directory.name)

    def tearDown(self):
        self.temporary_directory.cleanup()

    def test_pptx_converts_slides_title_text_and_table_with_locators(self):
        source = self.root / "产品介绍.pptx"
        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        slide.shapes.title.text = "产品概览"
        slide.placeholders[1].text = "第一段\n第二段"
        table_shape = slide.shapes.add_table(
            2, 2, Inches(1), Inches(4), Inches(5), Inches(1)
        )
        table_shape.table.cell(0, 0).text = "指标"
        table_shape.table.cell(0, 1).text = "数值"
        table_shape.table.cell(1, 0).text = "用户"
        table_shape.table.cell(1, 1).text = "100"
        presentation.save(source)

        first = PptxConverter().convert(source, self.root / "first")
        second = PptxConverter().convert(source, self.root / "second")
        markdown = strip_source_unit_markers(
            first.document.path.read_text(encoding="utf-8")
        )
        units = read_source_map(first.source_map.path)

        self.assertLess(markdown.index("# Slide 1"), markdown.index("## 产品概览"))
        self.assertLess(markdown.index("## 产品概览"), markdown.index("第一段"))
        self.assertLess(
            markdown.index("第一段"), markdown.index("| Column 1 | Column 2 |")
        )
        self.assertEqual(
            [unit["source"]["kind"] for unit in units],
            ["slide", "title", "text", "table"],
        )
        self.assertEqual(units[1]["source"]["slide"], 1)
        self.assertGreaterEqual(units[1]["source"]["shape"], 1)
        self.assertEqual(units[-1]["source"]["table"], 1)
        self.assertEqual(first.document.sha256, second.document.sha256)
        self.assertEqual(first.source_map.sha256, second.source_map.sha256)
        self.assertEqual(
            first.manifest_artifact.sha256, second.manifest_artifact.sha256
        )
        self.assertEqual(first.manifest.document_kind, "pptx")
        self.assertEqual(first.manifest.conversion_profile, "pptx-v1")

    def test_pptx_ignores_images_and_reports_warning(self):
        source = self.root / "图片.pptx"
        image = self.root / "pixel.png"
        image.write_bytes(
            base64.b64decode(
                "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR42mNk+A8AAQUBAScY42YAAAAASUVORK5CYII="
            )
        )
        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[5])
        slide.shapes.title.text = "有图片的页面"
        slide.shapes.add_picture(str(image), Inches(1), Inches(2))
        presentation.save(source)

        result = PptxConverter().convert(source, self.root / "derived")

        self.assertIn("PPTX_IMAGES_IGNORED:1", result.manifest.warnings)
        self.assertNotIn("pixel.png", result.document.path.read_text(encoding="utf-8"))

    def test_xlsx_uses_data_only_row_blocks_and_warns_for_uncached_formula(self):
        source = self.root / "监控记录.xlsx"
        workbook = openpyxl.Workbook()
        sheet = workbook.active
        sheet.title = "服务器监控"
        sheet.append(["主机", "负载"])
        sheet.append(["node-1", 0.5])
        sheet.append(["node-2", "=B2*2"])
        workbook.save(source)
        workbook.close()

        converter = XlsxConverter(spreadsheet_limits=SpreadsheetLimits(rows_per_unit=2))
        result = converter.convert(source, self.root / "derived")
        markdown = strip_source_unit_markers(
            result.document.path.read_text(encoding="utf-8")
        )
        units = read_source_map(result.source_map.path)

        self.assertIn("# Sheet: 服务器监控", markdown)
        self.assertIn("| 2 | node-1 | 0.5 |", markdown)
        self.assertIn("| 3 | node-2 |", markdown)
        self.assertNotIn("=B2*2", markdown)
        self.assertIn("XLSX_FORMULAS_WITHOUT_CACHED_VALUES:1", result.manifest.warnings)
        table_units = [unit for unit in units if unit["source"]["kind"] == "table"]
        self.assertEqual(
            [
                (unit["source"]["row_start"], unit["source"]["row_end"])
                for unit in table_units
            ],
            [(1, 2), (3, 3)],
        )
        self.assertEqual(table_units[0]["source"]["sheet"], "服务器监控")
        self.assertEqual(result.manifest.source_encoding, "binary")
        self.assertEqual(result.manifest.document_kind, "xlsx")

    def test_csv_detects_local_encoding_and_delimiter_and_preserves_row_ranges(self):
        source = self.root / "人员.csv"
        source.write_bytes(
            "姓名;备注\r\n张三;包含|符号\r\n李四;正常\r\n".encode("gb18030")
        )
        converter = CsvConverter(spreadsheet_limits=SpreadsheetLimits(rows_per_unit=2))

        first = converter.convert(source, self.root / "first")
        second = converter.convert(source, self.root / "second")
        markdown = strip_source_unit_markers(
            first.document.path.read_text(encoding="utf-8")
        )
        units = read_source_map(first.source_map.path)

        self.assertIn("| 2 | 张三 | 包含\\|符号 |", markdown)
        # The delimiter is provenance, not a defect: it must not push the file
        # into completed_with_warnings.
        self.assertIn("CSV_DELIMITER:semicolon", first.manifest.notes)
        self.assertEqual(first.manifest.warnings, ())
        self.assertNotEqual(first.manifest.source_encoding, "utf-8")
        self.assertEqual(first.manifest.document_kind, "csv")
        self.assertEqual(
            [
                (unit["source"]["row_start"], unit["source"]["row_end"])
                for unit in units
            ],
            [(1, 2), (3, 3)],
        )
        self.assertEqual(first.document.sha256, second.document.sha256)
        self.assertEqual(first.source_map.sha256, second.source_map.sha256)
        self.assertEqual(
            first.manifest_artifact.sha256, second.manifest_artifact.sha256
        )

    def test_macro_extensions_are_rejected_before_parsing(self):
        pptm = self.root / "macro.pptm"
        pptm.write_bytes(b"not parsed")
        with self.assertRaises(DocumentConversionError) as context:
            PptxConverter().convert(pptm, self.root / "pptm-derived")
        self.assertEqual(context.exception.code, "MACRO_ENABLED_DOCUMENT_UNSUPPORTED")

        xlsm = self.root / "macro.xlsm"
        xlsm.write_bytes(b"not parsed")
        with self.assertRaises(DocumentConversionError) as context:
            XlsxConverter().convert(xlsm, self.root / "xlsm-derived")
        self.assertEqual(context.exception.code, "MACRO_ENABLED_DOCUMENT_UNSUPPORTED")

    def test_invalid_and_suspicious_ooxml_packages_have_stable_codes(self):
        invalid = self.root / "invalid.xlsx"
        with zipfile.ZipFile(invalid, "w") as archive:
            archive.writestr("[Content_Types].xml", "<Types />")
        with self.assertRaises(DocumentConversionError) as context:
            XlsxConverter().convert(invalid, self.root / "invalid-derived")
        self.assertEqual(context.exception.code, "OOXML_INVALID_STRUCTURE")

        valid = self.root / "valid.xlsx"
        workbook = openpyxl.Workbook()
        workbook.active["A1"] = "content"
        workbook.save(valid)
        workbook.close()
        with self.assertRaises(DocumentConversionError) as context:
            XlsxConverter(zip_limits=OoxmlZipLimits(max_entries=2)).convert(
                valid,
                self.root / "entries-derived",
            )
        self.assertEqual(context.exception.code, "OOXML_ZIP_TOO_MANY_ENTRIES")

        compressed = self.root / "compressed.xlsx"
        with zipfile.ZipFile(
            compressed, "w", compression=zipfile.ZIP_DEFLATED
        ) as archive:
            archive.writestr("[Content_Types].xml", "x" * 20_000)
            archive.writestr("_rels/.rels", "rels")
            archive.writestr("xl/workbook.xml", "workbook")
        with self.assertRaises(DocumentConversionError) as context:
            XlsxConverter(zip_limits=OoxmlZipLimits(max_compression_ratio=2)).convert(
                compressed,
                self.root / "ratio-derived",
            )
        self.assertEqual(context.exception.code, "OOXML_ZIP_COMPRESSION_RATIO_EXCEEDED")

        unsafe_xml = self.root / "unsafe.xlsx"
        with zipfile.ZipFile(unsafe_xml, "w") as archive:
            archive.writestr(
                "[Content_Types].xml",
                '<!DOCTYPE x [<!ENTITY payload "unsafe">]><Types>&payload;</Types>',
            )
            archive.writestr("_rels/.rels", "<Relationships />")
            archive.writestr("xl/workbook.xml", "<workbook />")
        with self.assertRaises(DocumentConversionError) as context:
            XlsxConverter().convert(unsafe_xml, self.root / "unsafe-derived")
        self.assertEqual(context.exception.code, "OOXML_UNSAFE_XML")

    def test_spreadsheet_and_presentation_limits_have_stable_codes(self):
        xlsx = self.root / "rows.xlsx"
        workbook = openpyxl.Workbook()
        workbook.active.append(["one"])
        workbook.active.append(["two"])
        workbook.save(xlsx)
        workbook.close()
        with self.assertRaises(DocumentConversionError) as context:
            XlsxConverter(
                spreadsheet_limits=SpreadsheetLimits(max_rows_per_sheet=1)
            ).convert(
                xlsx,
                self.root / "rows-derived",
            )
        self.assertEqual(context.exception.code, "XLSX_TOO_MANY_ROWS")

        pptx = self.root / "slides.pptx"
        presentation = Presentation()
        presentation.slides.add_slide(
            presentation.slide_layouts[5]
        ).shapes.title.text = "one"
        presentation.slides.add_slide(
            presentation.slide_layouts[5]
        ).shapes.title.text = "two"
        presentation.save(pptx)
        with self.assertRaises(DocumentConversionError) as context:
            PptxConverter(presentation_limits=PresentationLimits(max_slides=1)).convert(
                pptx,
                self.root / "slides-derived",
            )
        self.assertEqual(context.exception.code, "PPTX_TOO_MANY_SLIDES")


if __name__ == "__main__":
    unittest.main()
