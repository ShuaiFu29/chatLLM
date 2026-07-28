from __future__ import annotations

import zipfile
from dataclasses import dataclass
from pathlib import Path
from xml.etree import ElementTree

from converted_document import ConversionResult, DocumentConversionError
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from converters.base import ConversionLimits, DocumentConverter, LocatedTextBlock
from converters.spreadsheet import (
    OoxmlZipLimits,
    _escape_heading,
    _escape_table_cell,
    inspect_ooxml_package,
)


@dataclass(frozen=True)
class PresentationLimits:
    max_slides: int = 500
    max_shapes_per_slide: int = 500
    max_tables_per_slide: int = 100
    max_table_rows: int = 1_000
    max_table_columns: int = 100
    max_text_chars: int = 65_536

    def __post_init__(self) -> None:
        if (
            self.max_slides <= 0
            or self.max_shapes_per_slide <= 0
            or self.max_tables_per_slide <= 0
            or self.max_table_rows <= 0
            or self.max_table_columns <= 0
            or self.max_text_chars <= 0
        ):
            raise ValueError("presentation limits must be positive")


class PptxConverter(DocumentConverter):
    document_kind = "pptx"
    conversion_profile = "pptx-v1"
    converter_name = "pptx-local"
    source_type = "pptx"

    def __init__(
        self,
        limits: ConversionLimits | None = None,
        *,
        presentation_limits: PresentationLimits | None = None,
        zip_limits: OoxmlZipLimits | None = None,
    ):
        super().__init__(limits)
        self.presentation_limits = presentation_limits or PresentationLimits()
        self.zip_limits = zip_limits or OoxmlZipLimits()

    def convert(
        self, source_file: str | Path, output_dir: str | Path
    ) -> ConversionResult:
        inspection = self._inspect_source(source_file)
        if inspection.path.suffix.lower() in {".pptm", ".potm", ".ppsm"}:
            raise DocumentConversionError(
                "MACRO_ENABLED_DOCUMENT_UNSUPPORTED",
                "macro-enabled presentation formats are not supported",
            )

        package = inspect_ooxml_package(
            inspection.path,
            self.zip_limits,
            required_entries=(
                "[Content_Types].xml",
                "_rels/.rels",
                "ppt/presentation.xml",
            ),
        )
        if package.has_macros:
            raise DocumentConversionError(
                "MACRO_ENABLED_DOCUMENT_UNSUPPORTED",
                "macro-enabled presentation packages are not supported",
            )

        warnings: list[str] = []
        if _has_external_relationships(inspection.path, package.names):
            warnings.append("PPTX_EXTERNAL_LINKS_IGNORED")

        try:
            presentation = Presentation(inspection.path)
        except (OSError, ValueError, KeyError, zipfile.BadZipFile) as error:
            raise DocumentConversionError(
                "PPTX_PARSE_FAILED", "presentation package could not be parsed"
            ) from error

        if len(presentation.slides) > self.presentation_limits.max_slides:
            raise DocumentConversionError(
                "PPTX_TOO_MANY_SLIDES",
                "presentation exceeds the configured slide limit",
            )

        units: list[LocatedTextBlock] = []
        has_content = False
        ignored_images = 0
        try:
            for slide_number, slide in enumerate(presentation.slides, start=1):
                shapes = list(slide.shapes)
                if len(shapes) > self.presentation_limits.max_shapes_per_slide:
                    raise DocumentConversionError(
                        "PPTX_TOO_MANY_SHAPES",
                        "presentation slide exceeds the configured shape limit",
                    )

                units.append(
                    LocatedTextBlock(
                        f"# Slide {slide_number}\n",
                        {"type": "pptx", "kind": "slide", "slide": slide_number},
                    )
                )
                title_shape = slide.shapes.title
                title_shape_id = (
                    title_shape.shape_id if title_shape is not None else None
                )

                if title_shape is not None:
                    title_text = _shape_text(
                        title_shape, self.presentation_limits.max_text_chars
                    )
                    if title_text:
                        has_content = True
                        shape_index = _shape_index(shapes, title_shape_id)
                        units.append(
                            LocatedTextBlock(
                                f"## {_escape_heading(title_text.replace(chr(10), ' '))}\n",
                                {
                                    "type": "pptx",
                                    "kind": "title",
                                    "slide": slide_number,
                                    "shape": shape_index,
                                    "shape_id": title_shape_id,
                                },
                            )
                        )

                # Text shapes precede tables while retaining their order within each category.
                for shape_index, shape in enumerate(shapes, start=1):
                    if shape.shape_id == title_shape_id or getattr(
                        shape, "has_table", False
                    ):
                        continue
                    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        ignored_images += 1
                        continue
                    text = _shape_text(shape, self.presentation_limits.max_text_chars)
                    if not text:
                        continue
                    has_content = True
                    units.append(
                        LocatedTextBlock(
                            text + "\n",
                            {
                                "type": "pptx",
                                "kind": "text",
                                "slide": slide_number,
                                "shape": shape_index,
                                "shape_id": shape.shape_id,
                            },
                        )
                    )

                table_index = 0
                for shape_index, shape in enumerate(shapes, start=1):
                    if not getattr(shape, "has_table", False):
                        continue
                    table_index += 1
                    if table_index > self.presentation_limits.max_tables_per_slide:
                        raise DocumentConversionError(
                            "PPTX_TOO_MANY_TABLES",
                            "presentation slide exceeds the configured table limit",
                        )
                    table_markdown, table_has_content = _table_markdown(
                        shape.table,
                        self.presentation_limits,
                    )
                    if not table_has_content:
                        continue
                    has_content = True
                    units.append(
                        LocatedTextBlock(
                            table_markdown,
                            {
                                "type": "pptx",
                                "kind": "table",
                                "slide": slide_number,
                                "shape": shape_index,
                                "shape_id": shape.shape_id,
                                "table": table_index,
                            },
                        )
                    )
        except DocumentConversionError:
            raise
        except (AttributeError, KeyError, TypeError, ValueError) as error:
            raise DocumentConversionError(
                "PPTX_PARSE_FAILED", "presentation content could not be parsed"
            ) from error

        if not has_content:
            raise DocumentConversionError(
                "EMPTY_DOCUMENT", "presentation has no embeddable text"
            )
        if ignored_images:
            warnings.append(f"PPTX_IMAGES_IGNORED:{ignored_images}")

        return self._convert_located_blocks(
            inspection,
            output_dir,
            units,
            "binary",
            tuple(warnings),
        )


def _shape_text(shape, max_chars: int) -> str:
    if not getattr(shape, "has_text_frame", False):
        return ""
    paragraphs: list[str] = []
    for paragraph in shape.text_frame.paragraphs:
        text = _normalize_pptx_text(paragraph.text, max_chars)
        if text:
            paragraphs.append(text)
    combined = "\n\n".join(paragraphs)
    if len(combined) > max_chars:
        raise DocumentConversionError(
            "PPTX_TEXT_TOO_LARGE",
            "presentation text shape exceeds the configured character limit",
        )
    return combined


def _normalize_pptx_text(value: str, max_chars: int) -> str:
    normalized = value.replace("\r\n", "\n").replace("\r", "\n").strip()
    normalized = normalized.replace("<!-- source-unit:", "&lt;!-- source-unit:")
    if len(normalized) > max_chars:
        raise DocumentConversionError(
            "PPTX_TEXT_TOO_LARGE",
            "presentation text exceeds the configured character limit",
        )
    return normalized


def _table_markdown(table, limits: PresentationLimits) -> tuple[str, bool]:
    row_count = len(table.rows)
    column_count = len(table.columns)
    if row_count > limits.max_table_rows:
        raise DocumentConversionError(
            "PPTX_TABLE_TOO_MANY_ROWS",
            "presentation table exceeds the configured row limit",
        )
    if column_count > limits.max_table_columns:
        raise DocumentConversionError(
            "PPTX_TABLE_TOO_MANY_COLUMNS",
            "presentation table exceeds the configured column limit",
        )
    if not row_count or not column_count:
        return "", False

    headers = [f"Column {index}" for index in range(1, column_count + 1)]
    lines = [
        "| " + " | ".join(headers) + " |",
        "| " + " | ".join("---" for _ in headers) + " |",
    ]
    has_content = False
    for row in table.rows:
        values: list[str] = []
        for cell in row.cells:
            value = _normalize_pptx_text(cell.text, limits.max_text_chars)
            has_content = has_content or bool(value)
            values.append(_escape_table_cell(value))
        lines.append("| " + " | ".join(values) + " |")
    return "\n".join(lines) + "\n", has_content


def _shape_index(shapes: list, shape_id: int) -> int:
    for index, shape in enumerate(shapes, start=1):
        if shape.shape_id == shape_id:
            return index
    raise DocumentConversionError(
        "PPTX_PARSE_FAILED", "title shape locator could not be resolved"
    )


def _has_external_relationships(path: Path, names: tuple[str, ...]) -> bool:
    relationship_names = sorted(
        name for name in names if name.casefold().endswith(".rels")
    )
    try:
        with zipfile.ZipFile(path) as archive:
            for name in relationship_names:
                with archive.open(name) as stream:
                    root = ElementTree.parse(stream).getroot()
                for relationship in root:
                    if (
                        relationship.attrib.get("TargetMode", "").casefold()
                        == "external"
                    ):
                        return True
    except (OSError, zipfile.BadZipFile, ElementTree.ParseError) as error:
        raise DocumentConversionError(
            "PPTX_PARSE_FAILED", "presentation relationships could not be parsed"
        ) from error
    return False
